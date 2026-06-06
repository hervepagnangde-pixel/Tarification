"""
Atlantic Re IA — PDF & PPTX generation module
ReportLab pour les rapports PDF, PptxGenJS (Node.js) pour le PowerPoint.
"""
import io as _io_db
import streamlit as st
import os, json, subprocess, tempfile
from datetime import datetime
import numpy as np
import pandas as pd
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
    TableStyle, PageBreak, HRFlowable)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

# ════════════════════════════════════════════
# GENERATION PDF PROFESSIONNEL
# ════════════════════════════════════════════
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
    TableStyle, PageBreak, HRFlowable)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
import io as _io_db
import streamlit as st

_VERT  = colors.HexColor("#2d8a4e")
_NOIR  = colors.HexColor("#1a1a1a")
_GRIS  = colors.HexColor("#f4f6f4")
_GRIS2 = colors.HexColor("#888888")

def _pdf_styles():
    s = getSampleStyleSheet()
    for nm, kwargs in [
        ("AR_Title",  dict(fontName="Helvetica-Bold", fontSize=26, textColor=_NOIR,   spaceAfter=6,  alignment=TA_CENTER)),
        ("AR_Sub",    dict(fontName="Helvetica",      fontSize=12, textColor=_GRIS2,  spaceAfter=4,  alignment=TA_CENTER)),
        ("AR_H1",     dict(fontName="Helvetica-Bold", fontSize=14, textColor=_VERT,   spaceBefore=14,spaceAfter=6)),
        ("AR_H2",     dict(fontName="Helvetica-Bold", fontSize=10, textColor=_NOIR,   spaceBefore=8, spaceAfter=3)),
        ("AR_Body",   dict(fontName="Helvetica",      fontSize=9,  leading=14,        textColor=_NOIR,spaceAfter=4)),
        ("AR_Caption",dict(fontName="Helvetica-Oblique",fontSize=8,textColor=_GRIS2,  spaceAfter=8)),
        ("AR_Cell",   dict(fontName="Helvetica",      fontSize=8,  leading=11,        textColor=_NOIR)),
        ("AR_CellB",  dict(fontName="Helvetica-Bold", fontSize=8,  leading=11,        textColor=_NOIR)),
    ]:
        if nm not in s.byName:
            kwargs.pop("alignment", None)
            s.add(ParagraphStyle(nm, parent=s["Normal"], **kwargs))
    return s

def _pdf_table_rl(data_rows, col_widths=None):
    S = _pdf_styles()
    rows = []
    for i, row in enumerate(data_rows):
        r = []
        for cell in row:
            sty = S["AR_CellB"] if i == 0 else S["AR_Cell"]
            r.append(Paragraph(str(cell), sty))
        rows.append(r)
    t = Table(rows, colWidths=col_widths, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),  (-1,0),  _NOIR),
        ("TEXTCOLOR",     (0,0),  (-1,0),  colors.white),
        ("ROWBACKGROUNDS",(0,1),  (-1,-1), [colors.white, _GRIS]),
        ("GRID",          (0,0),  (-1,-1), 0.3, colors.HexColor("#e0e0e0")),
        ("VALIGN",        (0,0),  (-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),  (-1,-1), 4),
        ("BOTTOMPADDING", (0,0),  (-1,-1), 4),
        ("LEFTPADDING",   (0,0),  (-1,-1), 5),
    ]))
    return t

def generer_pdf_rapport(user_email, gnpi_val, tranches,
                         resultats_bc, resultats_sim, taux_mkt_final,
                         df_rapport, prime_totale,
                         analyse_claude="", annee=2026):
    buf = _io_db.BytesIO()
    S   = _pdf_styles()
    W, H = A4
    m   = 1.8*cm
    doc = SimpleDocTemplate(buf, pagesize=A4,
          leftMargin=m, rightMargin=m, topMargin=1.5*cm, bottomMargin=1.5*cm,
          title=f"Rapport Tarification Atlantic Re {annee}")
    story = []
    PW = W - 2*m

    # ── COUVERTURE ──
    story += [Spacer(1, 1.5*cm),
              Paragraph("ATLANTIC RE", ParagraphStyle("T", parent=S["Normal"],
                  fontName="Helvetica-Bold", fontSize=32, textColor=_NOIR, alignment=1,
                  spaceAfter=6)),
              Paragraph(f"Rapport de Tarification {annee}", ParagraphStyle("S",
                  parent=S["Normal"], fontName="Helvetica", fontSize=14,
                  textColor=_GRIS2, alignment=1, spaceAfter=4)),
              Paragraph("Reassurance Non-Proportionnelle · Automobile · Maroc",
                  ParagraphStyle("S2", parent=S["Normal"], fontName="Helvetica",
                  fontSize=10, textColor=_GRIS2, alignment=1, spaceAfter=16)),
              HRFlowable(width=PW, thickness=2, color=_VERT),
              Spacer(1, 0.6*cm)]

    cv = [["GNPI", f"{gnpi_val:,.0f} MAD"],
          ["Prime totale", f"{prime_totale:,.0f} MAD"],
          ["Taux global", f"{prime_totale/gnpi_val:.4%}" if gnpi_val else "-"],
          ["Tranches", str(len(tranches))],
          ["Date", datetime.now().strftime("%d/%m/%Y %H:%M")],
          ["Prepare par", "Atlantic Re IA - Agent Actuariel"]]
    story.append(_pdf_table_rl([["Indicateur","Valeur"]]+cv, [PW*0.45, PW*0.55]))
    story.append(PageBreak())

    # ── 1. PROGRAMME ──
    story += [Paragraph("1. Programme de Reassurance", S["AR_H1"]),
              HRFlowable(width=PW, thickness=1, color=_VERT), Spacer(1, 0.3*cm)]
    ph = ["Tranche","Type","Priorite (MAD)","Portee (MAD)","Reconst.","Brokage","Frais","Marge"]
    pr = [[t.get("nom",""), t.get("type",""), f"{t.get('priorite',0):,.0f}",
           f"{t.get('portee',0):,.0f}",
           f"{t.get('nb_reconstitutions',1)}x{t.get('taux_reconstitution',100):.0f}%",
           f"{t.get('brokage',0):.0%}", f"{t.get('frais',0):.0%}", f"{t.get('marge',0):.0%}"]
          for t in tranches]
    story.append(_pdf_table_rl([ph]+pr,
        [PW*0.16,PW*0.13,PW*0.13,PW*0.12,PW*0.12,PW*0.10,PW*0.10,PW*0.10]))
    story.append(Spacer(1, 0.4*cm))

    # ── 2. BURNING COST ──
    if resultats_bc:
        story += [Paragraph("2. Burning Cost", S["AR_H1"]),
                  HRFlowable(width=PW, thickness=1, color=_VERT), Spacer(1, 0.3*cm),
                  Paragraph("Ck = min(max(S'k - D, 0), L) x coeff_stab | "
                             "R1: tau_risque = tau_pur + sigma_hist x 20% | "
                             "R2: tau_BC=0 si annees non-nulles < 3 | "
                             "tau_tech = tau_risque x (1-Rec) / (1-BK-FG-Marge-Retro)", S["AR_Caption"])]
        bh = ["Tranche","Type","Charge moy.","Rec","Taux pur","Taux risque","Taux tech.","Charg. majeurs"]
        br = [[r.get("tranche",""), r.get("type",""),
               f"{r.get('charge_moy', r.get('charge_moy_MAD',0)):,.0f}",
               f"{r.get('Rec',0):.4%}", f"{r.get('taux_pur',0):.4%}",
               f"{r.get('taux_risque',0):.4%}", f"{r.get('taux_technique',0):.4%}",
               f"{r.get('chargement_majeurs',0):.4%}"] for r in resultats_bc]
        story.append(_pdf_table_rl([bh]+br,
            [PW*0.18,PW*0.12,PW*0.14,PW*0.10,PW*0.11,PW*0.11,PW*0.12,PW*0.12]))
        story.append(Spacer(1, 0.4*cm))

    # ── 3. SIMULATION ──
    if resultats_sim:
        story += [Paragraph("3. Simulation Pareto / Poisson", S["AR_H1"]),
                  HRFlowable(width=PW, thickness=1, color=_VERT), Spacer(1, 0.3*cm)]
        sh = ["Tranche","Taux pur","Taux risque","Taux tech.","Charg. majeurs","Sans AAL","Sans AAD","Sans reconst."]
        sr = [[r.get("tranche",""), f"{r.get('taux_pur',0):.4%}", f"{r.get('taux_risque',0):.4%}",
               f"{r.get('taux_technique',0):.4%}", f"{r.get('chargement_majeurs',0):.4%}",
               f"{r.get('sans_aal',0):.4%}", f"{r.get('sans_aad',0):.4%}",
               f"{r.get('sans_rec',0):.4%}"] for r in resultats_sim]
        story.append(_pdf_table_rl([sh]+sr,
            [PW*0.17,PW*0.11,PW*0.12,PW*0.11,PW*0.11,PW*0.10,PW*0.10,PW*0.14]))
        story.append(Spacer(1, 0.4*cm))

    # ── 4. MARKET CURVE ──
    if taux_mkt_final:
        story += [Paragraph("4. Market Curve  ROL = a x x^(-b)", S["AR_H1"]),
                  HRFlowable(width=PW, thickness=1, color=_VERT), Spacer(1, 0.3*cm)]
        mh = ["Tranche","Type","x=(D+C/2)/GNPI","ROL","Taux pur","Taux tech.","Taux final"]
        mr = [[t.get("tranche",""), t.get("type",""), f"{t.get('x_norm',0):.5f}",
               f"{t.get('rol',0):.4%}", f"{t.get('taux_pur',0):.4%}",
               f"{t.get('taux_tech',0):.4%}", f"{t.get('taux',0):.4%}"]
              for t in taux_mkt_final]
        story.append(_pdf_table_rl([mh]+mr,
            [PW*0.17,PW*0.12,PW*0.16,PW*0.11,PW*0.11,PW*0.11,PW*0.12]))
        story.append(Spacer(1, 0.4*cm))

    # ── 5. SYNTHESE FINALE ──
    story += [PageBreak(),
              Paragraph("5. Synthese de Tarification", S["AR_H1"]),
              HRFlowable(width=PW, thickness=2, color=_VERT), Spacer(1, 0.3*cm)]

    if df_rapport is not None and not df_rapport.empty:
        cols = list(df_rapport.columns)
        drows = [cols] + [list(map(str, r)) for r in df_rapport.values.tolist()]
        cw = PW / len(cols)
        story.append(_pdf_table_rl(drows, [cw]*len(cols)))
        story.append(Spacer(1, 0.4*cm))

    story.append(_pdf_table_rl(
        [["Indicateur","Valeur"],
         ["Prime totale", f"{prime_totale:,.0f} MAD"],
         ["Taux global",  f"{prime_totale/gnpi_val:.4%}" if gnpi_val else "-"],
         ["GNPI",         f"{gnpi_val:,.0f} MAD"]],
        [PW*0.45, PW*0.55]))

    # ── 6. ANALYSE CLAUDE ──
    if analyse_claude and len(analyse_claude.strip()) > 10:
        story += [PageBreak(),
                  Paragraph("6. Analyse Claude - Recommandations", S["AR_H1"]),
                  HRFlowable(width=PW, thickness=1, color=_VERT), Spacer(1, 0.3*cm)]
        for line in analyse_claude.split("\n"):
            line = line.strip()
            if not line: story.append(Spacer(1, 0.15*cm)); continue
            line_c = line.replace("**","").replace("*","").replace("`","").replace("#","")
            if line.startswith("## "): story.append(Paragraph(line[3:], S["AR_H2"]))
            elif line.startswith("# "): story.append(Paragraph(line[2:], S["AR_H1"]))
            else: story.append(Paragraph(line_c, S["AR_Body"]))

    # ── FOOTER ──
    story += [Spacer(1, 0.8*cm),
              HRFlowable(width=PW, thickness=0.5, color=_GRIS2),
              Spacer(1, 0.2*cm),
              Paragraph(f"Document genere par Atlantic Re IA - Agent Actuariel Autonome | "
                        f"Atlantic Re {annee} | {datetime.now().strftime('%d/%m/%Y %H:%M')} | Confidentiel",
                        S["AR_Caption"])]

    doc.build(story)
    pdf_data = buf.getvalue()
    # ── Signature SHA-256 ────────────────────────────────────────
    sig = _pdf_signature_qr(pdf_data)
    return pdf_data


def _pdf_signature_qr(pdf_bytes):
    """Génère un hash SHA-256 du contenu PDF pour la traçabilité."""
    import hashlib
    h = hashlib.sha256(pdf_bytes).hexdigest()
    return h


def envoyer_webhook_notification(sujet, corps_texte, niveau="info"):
    """
    Envoie une notification via webhook Slack ou Microsoft Teams.
    Configurer dans Secrets : SLACK_WEBHOOK_URL ou TEAMS_WEBHOOK_URL
    """
    import urllib.request, urllib.error
    slack_url  = ""
    teams_url  = ""
    try:
        slack_url  = st.secrets.get("SLACK_WEBHOOK_URL", "")
        teams_url  = st.secrets.get("TEAMS_WEBHOOK_URL", "")
    except Exception:
        pass

    icone = {"info":"ℹ️","alerte":"⚠️","rapport_final":"📋","succes":"✅"}.get(niveau,"📊")
    resultats = []

    # ── Slack ──────────────────────────────────────────────────────
    if slack_url:
        payload = json.dumps({
            "text": f"{icone} *[Atlantic Re IA]* {sujet}",
            "blocks": [
                {"type":"section","text":{"type":"mrkdwn",
                    "text":f"{icone} *{sujet}*\n{corps_texte[:500]}"}},
                {"type":"context","elements":[
                    {"type":"mrkdwn","text":f"Atlantic Re IA · {datetime.now().strftime('%d/%m/%Y %H:%M')}"}]}
            ]
        }).encode()
        try:
            req = urllib.request.Request(slack_url, data=payload,
                headers={"Content-Type":"application/json"})
            urllib.request.urlopen(req, timeout=5)
            resultats.append(("Slack", True, "OK"))
        except Exception as e:
            resultats.append(("Slack", False, str(e)[:80]))

    # ── Teams ──────────────────────────────────────────────────────
    if teams_url:
        color_map = {"info":"0076D7","alerte":"FF8C00","rapport_final":"107C10","succes":"107C10"}
        payload = json.dumps({
            "@type":"MessageCard","@context":"http://schema.org/extensions",
            "themeColor": color_map.get(niveau,"0076D7"),
            "summary": sujet,
            "sections":[{"activityTitle": f"{icone} {sujet}",
                          "activitySubtitle": "Atlantic Re IA",
                          "text": corps_texte[:500]}]
        }).encode()
        try:
            req = urllib.request.Request(teams_url, data=payload,
                headers={"Content-Type":"application/json"})
            urllib.request.urlopen(req, timeout=5)
            resultats.append(("Teams", True, "OK"))
        except Exception as e:
            resultats.append(("Teams", False, str(e)[:80]))

    return resultats





def generer_pptx_rapport(gnpi_val, tranches, resultats_bc, resultats_sim,
                          taux_mkt_final, df_rapport, prime_totale, annee=2026):
    """
    Génère un rapport PPTX avec python-pptx (100% Python, fonctionne sur Streamlit Cloud).
    Palette : navy #0d2b3e + teal #00b5a5.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io as _io_p
    except ImportError:
        return None

    NAV  = RGBColor(0x0d, 0x2b, 0x3e)
    TEAL = RGBColor(0x00, 0xb5, 0xa5)
    WHT  = RGBColor(0xFF, 0xFF, 0xFF)
    GRY  = RGBColor(0xf2, 0xf8, 0xf7)
    DRK  = RGBColor(0x1a, 0x1a, 0x1a)

    taux_global = prime_totale / gnpi_val if gnpi_val else 0

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # Blank layout

    def add_slide():
        return prs.slides.add_slide(blank)

    def bg(slide, color):
        from pptx.oxml.ns import qn
        from lxml import etree
        bg_elem = slide.background
        fill = bg_elem.fill
        fill.solid()
        fill.fore_color.rgb = color

    def txbox(slide, text, left, top, width, height,
              size=18, bold=False, color=WHT, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return tb

    def rect(slide, left, top, width, height, color):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_table(slide, headers, rows, left, top, width, height):
        from pptx.util import Inches, Pt
        n_cols = len(headers)
        n_rows = len(rows) + 1
        tbl = slide.shapes.add_table(n_rows, n_cols,
            Inches(left), Inches(top), Inches(width), Inches(height)).table
        col_w = Inches(width / n_cols)
        for i in range(n_cols):
            tbl.columns[i].width = col_w
        # Header row
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAV
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(h)
            run.font.bold = True; run.font.size = Pt(9); run.font.color.rgb = WHT
        # Data rows
        for i, row in enumerate(rows):
            bg_color = GRY if i % 2 == 0 else WHT
            for j, val in enumerate(row):
                cell = tbl.cell(i+1, j)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg_color
                p = cell.text_frame.paragraphs[0]
                run = p.add_run(); run.text = str(val)
                run.font.size = Pt(8); run.font.color.rgb = DRK

    # ── Slide 1 : Couverture ──────────────────────────────────────────
    s1 = add_slide(); bg(s1, NAV)
    rect(s1, 0, 5.8, 13.33, 0.5, TEAL)
    txbox(s1, "ATLANTIC RE", 0.5, 0.8, 12, 1.5, size=48, bold=True, color=WHT)
    txbox(s1, f"Rapport de Tarification {annee}", 0.5, 2.4, 12, 0.8, size=22, color=TEAL)
    txbox(s1, "Réassurance Non-Proportionnelle · Automobile · Maroc",
          0.5, 3.1, 12, 0.6, size=13, color=WHT)
    txbox(s1, f"GNPI : {gnpi_val:,.0f} MAD  |  Prime totale : {prime_totale:,.0f} MAD  |  Taux global : {taux_global:.4%}",
          0.5, 5.85, 12, 0.4, size=11, color=NAV)

    # ── Slide 2 : Programme ───────────────────────────────────────────
    s2 = add_slide(); bg(s2, GRY)
    rect(s2, 0, 0, 13.33, 0.85, NAV)
    txbox(s2, "Programme de Réassurance", 0.4, 0.05, 12, 0.75, size=22, bold=True, color=WHT)
    rect(s2, 0, 0.85, 13.33, 0.05, TEAL)
    if tranches:
        hdrs = ["Tranche","Type","Priorité","Portée","Reconstitutions"]
        rows = [[t.get("nom",""), t.get("type",""),
                 f"{t.get('priorite',0)/1e6:.1f}M MAD",
                 f"{t.get('portee',0)/1e6:.1f}M MAD",
                 f"{t.get('nb_reconstitutions',1)}x{t.get('taux_reconstitution',100):.0f}%"]
                for t in tranches]
        add_table(s2, hdrs, rows, 0.4, 1.1, 12.5, min(0.45*len(rows)+0.5, 5.5))

    # ── Slide 3 : Burning Cost ────────────────────────────────────────
    s3 = add_slide(); bg(s3, GRY)
    rect(s3, 0, 0, 13.33, 0.85, NAV)
    txbox(s3, "Burning Cost", 0.4, 0.05, 12, 0.75, size=22, bold=True, color=WHT)
    rect(s3, 0, 0.85, 13.33, 0.05, TEAL)
    txbox(s3, "Méthode de référence · R1 : τ_risque = τ_pur + σ × 20%",
          0.4, 0.95, 12, 0.35, size=10, italic=True, color=DRK)
    if resultats_bc:
        hdrs = ["Tranche","τ pur","τ risque","τ technique","Années non nulles"]
        rows = [[r.get("tranche",""), f"{r.get('taux_pur',0):.4%}",
                 f"{r.get('taux_risque',0):.4%}", f"{r.get('taux_technique',0):.4%}",
                 str(r.get('n_ann_nonzero',0))] for r in resultats_bc]
        add_table(s3, hdrs, rows, 0.4, 1.4, 12.5, min(0.45*len(rows)+0.5, 5.5))

    # ── Slide 4 : Simulation ──────────────────────────────────────────
    s4 = add_slide(); bg(s4, GRY)
    rect(s4, 0, 0, 13.33, 0.85, NAV)
    txbox(s4, "Simulation Pareto / Poisson", 0.4, 0.05, 12, 0.75, size=22, bold=True, color=WHT)
    rect(s4, 0, 0.85, 13.33, 0.05, TEAL)
    if resultats_sim:
        hdrs = ["Tranche","τ pur","τ risque","τ technique","Sans AAL","Sans AAD"]
        rows = [[r.get("tranche",""), f"{r.get('taux_pur',0):.4%}",
                 f"{r.get('taux_risque',0):.4%}", f"{r.get('taux_technique',0):.4%}",
                 f"{r.get('sans_aal',0):.4%}", f"{r.get('sans_aad',0):.4%}"]
                for r in resultats_sim]
        add_table(s4, hdrs, rows, 0.4, 1.1, 12.5, min(0.45*len(rows)+0.5, 5.5))

    # ── Slide 5 : Synthèse ────────────────────────────────────────────
    s5 = add_slide(); bg(s5, GRY)
    rect(s5, 0, 0, 13.33, 0.85, NAV)
    txbox(s5, "Synthèse de Tarification", 0.4, 0.05, 12, 0.75, size=22, bold=True, color=WHT)
    rect(s5, 0, 0.85, 13.33, 0.05, TEAL)
    if df_rapport is not None and not df_rapport.empty:
        cols = list(df_rapport.columns)
        rows = [list(map(str, r)) for r in df_rapport.values.tolist()]
        add_table(s5, cols, rows, 0.4, 1.1, 12.5, min(0.45*len(rows)+0.5, 5.8))

    # ── Slide 6 : Conclusion ──────────────────────────────────────────
    s6 = add_slide(); bg(s6, NAV)
    rect(s6, 0, 0, 13.33, 0.06, TEAL)
    txbox(s6, "Conclusion & Recommandations", 0.5, 0.4, 9, 1.0, size=28, bold=True, color=WHT)
    txbox(s6, "Prime totale", 0.8, 1.6, 4, 0.4, size=13, bold=True, color=TEAL)
    txbox(s6, f"{prime_totale:,.0f} MAD", 0.8, 2.0, 4, 0.7, size=24, bold=True, color=WHT)
    txbox(s6, "Taux global", 0.8, 2.8, 4, 0.4, size=13, bold=True, color=TEAL)
    txbox(s6, f"{taux_global:.4%}", 0.8, 3.2, 4, 0.7, size=24, bold=True, color=WHT)
    txbox(s6, (f"Sélection : max(BC, Sim) travaillante\n"
               f"max(Sim, Marché) cat / non-travaillante\n\n"
               f"BC = méthode de référence\n"
               f"Simulation = validation & prudence\n"
               f"Généré par Atlantic Re IA · {datetime.now().strftime('%d/%m/%Y %H:%M')}"),
          5.5, 1.8, 7.0, 4.0, size=12, color=RGBColor(0xc8, 0xdc, 0xe6))

    buf = _io_p.BytesIO()
    prs.save(buf)
    return buf.getvalue()
